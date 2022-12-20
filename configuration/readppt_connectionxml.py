
import pandas as pd
from pptx import Presentation
#######################################################
modeltype_all2=['pv', 'wind', 'load', 'battery', 'h2storage']
file="example.pptx"
config_model=[]
prs = Presentation(file)
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if shape.text.lower() in modeltype_all2:
                config_model.append(shape.text)
###########################################################
#config_model=['PV','Wind','Load','Battery','H2storage']

###########################################################
config_model=pd.DataFrame(config_model,columns=['model'])
config_model.dropna()
config_model.reset_index()
config_data=config_model.groupby(config_model.columns.tolist(),as_index=False).size()
connection=[]
for model_i in config_data.iterrows():
    model_number = int(model_i[1]['size'])
    model_type = model_i[1]['model']
    for number in range(model_number):
        if model_type[0:2].lower()=='wi':
            connection.append(['wind['+ str(number) +']', 'ctrl', 'wind_gen'])
            connection.append(['wind['+ str(number) +']', 'monitor', 'wind_gen'],)
        elif model_type[0:2].lower()=='pv':
            connection.append(['pv[' + str(number) + ']', 'ctrl', 'pv_gen'])
            connection.append(['pv[' + str(number) + ']', 'monitor', 'pv_gen'],)
        elif model_type[0:2].lower()=='lo':
            connection.append(['load[' + str(number) + ']', 'ctrl', 'load_dem'])
            connection.append(['load[' + str(number) + ']', 'monitor', 'load_dem'],)
        elif model_type[0:2].lower()=='ba':
            connection.append(['ctrl', 'battery', 'flow2b', "async_requests=True"])
            connection.append(['battery', 'monitor', 'p_out'])
            connection.append(['battery', 'monitor', 'p_in'])
            connection.append(['battery', 'monitor',  'mod'])
            connection.append(['battery', 'monitor',  'flag'])
            connection.append(['battery', 'monitor',  'soc'])
            connection.append(['ctrl', 'monitor', 'flow2b'])
        elif model_type[0:2].lower()=='h2':
            connection.append(['ctrl', 'electrolyser', 'flow2e'])
            connection.append(['electrolyser', 'h2storage', 'h2_gen', 'h2_in'])
            connection.append(['h2storage', 'fuelcell', 'h2_given', 'h2_consume'])
            connection.append(['ctrl', 'h2storage', 'h2_out', "async_requests=True"])
            connection.append(['ctrl', 'monitor', 'flow2e'])
            connection.append(['ctrl', 'monitor', 'dump'])
            connection.append(['electrolyser', 'monitor', 'h2_gen'])
            connection.append(['h2storage', 'monitor', 'h2_stored'])
            connection.append(['h2storage', 'monitor', 'h2_given'])
            connection.append(['h2storage', 'monitor', 'h2_stored'])
            connection.append(['h2storage', 'monitor', 'h2_soc'])
            connection.append(['fuelcell', 'monitor', 'fc_gen'])
try:
    df=pd.DataFrame(connection, columns=['send','receive','message','more'])
except ValueError:
    df=pd.DataFrame(connection, columns=['send','receive','message'])
data=df.to_xml()

with open('connection.xml','w') as file:
    file.write(data)




sim_config=[['Wind' ,'python','Models.Wind.wind_mosaik:WindSim'],
            ['PV' ,'python', 'Models.PV.pv_mosaik:PvAdapter'],
            ['Load','python', 'Models.Load.load_mosaik:loadSim'],
            ['Collector','python', 'Models.collector:Collector'],
            ['CSVB','python', 'mosaik_csv:CSV'],
            ['Controller','python', 'Models.Controller.controller_mosaik:controlSim'],
            ['Battery','python', 'Models.Battery.battery_mosaik:BatteryholdSim'],
            #['H2storage','connect', '192.168.0.2:5123'],
            #['PV','connect', '192.168.0.1:5123'],
            ['Electrolyser','python', 'Models.Electrolyser.electrolyser_mosaik:ElectrolyserSim'],
            ['H2storage','python', 'Models.H2storage.h2storage_mosaik:compressedhydrogen'],
            ['Fuelcell','python', 'Models.Fuelcell.fuelcell_mosaik:FuelCellSim'],
            ]
sim_config_df=pd.DataFrame(sim_config, columns=['model','method','location'])
sim_config_data=sim_config_df.to_xml()

with open('config.xml','w') as file:
    file.write(sim_config_data)

